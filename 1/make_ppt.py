"""
2025年度工作汇报 PPT Generator
iOS-style Flat Design · Frontend Developer · CATL Energy Storage
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ── Color Palette (iOS style) ──
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLUE       = RGBColor(0x00, 0x7A, 0xFF)
DARK       = RGBColor(0x1D, 0x1D, 0x1F)
GRAY       = RGBColor(0x86, 0x86, 0x8B)
LIGHT_BG   = RGBColor(0xF5, 0xF5, 0xF7)
GREEN      = RGBColor(0x34, 0xC7, 0x59)
ORANGE     = RGBColor(0xFF, 0x95, 0x00)
RED_SOFT   = RGBColor(0xFF, 0x3B, 0x30)
BORDER     = RGBColor(0xE5, 0xE5, 0xEA)
CARD_BG    = RGBColor(0xFA, 0xFA, 0xFA)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN  = Inches(0.9)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
BLANK_LAYOUT = prs.slide_layouts[6]


# ═══════════════════════════════════════════════
#  HELPER FUNCTIONS
# ═══════════════════════════════════════════════

def add_top_accent(slide, color=BLUE, width=None, height=Inches(0.06)):
    """Thin colored bar at the very top of a content slide."""
    w = width or SLIDE_W
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, w, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def add_textbox(slide, left, top, width, height, text, font_size=Pt(16),
                color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    """Add a simple text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    # Set East Asian font
    for run in p.runs:
        rPr = run._r.get_or_add_rPr()
        rPr.set(qn('a:eaTypeface'), font_name)
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines, font_size=Pt(14),
                          color=DARK, line_spacing=Pt(24), font_name='Microsoft YaHei'):
    """Add a text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = font_size
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = line_spacing
    return txBox


def add_card(slide, left, top, width, height, fill_color=WHITE, border_color=BORDER, radius=None):
    """Add a rounded rectangle card."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(0.5)
    if radius:
        shape.adjustments[0] = radius
    return shape


def add_card_with_text(slide, left, top, width, height, title, body_lines,
                       title_size=Pt(18), body_size=Pt(13), body_color=GRAY):
    """Card with title + body text."""
    add_card(slide, left, top, width, height)
    pad = Inches(0.25)
    add_textbox(slide, left + pad, top + pad + Inches(0.05),
                width - 2*pad, Inches(0.4), title, font_size=title_size, bold=True)
    add_multiline_textbox(slide, left + pad, top + pad + Inches(0.45),
                          width - 2*pad, height - Inches(0.6),
                          body_lines, font_size=body_size, color=body_color,
                          line_spacing=Pt(20))


def add_metric_card(slide, left, top, width, height, number, label, num_color=BLUE):
    """Big number + label card."""
    add_card(slide, left, top, width, height)
    add_textbox(slide, left, top + Inches(0.15), width, Inches(0.7),
                number, font_size=Pt(28), color=num_color, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left, top + Inches(0.75), width, Inches(0.5),
                label, font_size=Pt(12), color=GRAY, alignment=PP_ALIGN.CENTER)


def add_bullet_circle(slide, left, top, size, color=BLUE, fill=True):
    """Small circle for bullet/number decoration."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
    else:
        shape.fill.background()
    shape.line.color.rgb = color
    shape.line.width = Pt(1.5)
    return shape


def add_page_number(slide, num, total=12):
    """Page number at bottom right."""
    add_textbox(slide, SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.5),
                Inches(0.8), Inches(0.35), f"{num} / {total}",
                font_size=Pt(10), color=GRAY, alignment=PP_ALIGN.RIGHT)


def add_slide_title(slide, text, subtitle=None):
    """Standard slide title with optional subtitle."""
    add_top_accent(slide)
    add_textbox(slide, MARGIN, Inches(0.3), Inches(11), Inches(0.7),
                text, font_size=Pt(30), bold=True, color=DARK)
    if subtitle:
        add_textbox(slide, MARGIN, Inches(0.9), Inches(11), Inches(0.4),
                    subtitle, font_size=Pt(14), color=GRAY)


# ═══════════════════════════════════════════════
#  SLIDE 1 - COVER
# ═══════════════════════════════════════════════
def make_cover():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    # Full blue accent block on left
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.15), SLIDE_H)
    accent.fill.solid()
    accent.fill.fore_color.rgb = BLUE
    accent.line.fill.background()

    # Large circle decoration (subtle)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     SLIDE_W - Inches(5), Inches(-1.5),
                                     Inches(7), Inches(7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xFE)
    circle.line.fill.background()

    # Title
    add_textbox(slide, MARGIN, Inches(1.5), Inches(10), Inches(1.2),
                "2025年度工作汇报", font_size=Pt(44), bold=True, color=DARK)
    # Subtitle
    add_textbox(slide, MARGIN, Inches(2.8), Inches(10), Inches(0.6),
                "前端开发  ·  宁德时代储能板块", font_size=Pt(22), color=BLUE)
    # Horizontal line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, Inches(3.6),
                                   Inches(2.5), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()
    # Name / Date
    add_textbox(slide, MARGIN, Inches(4.0), Inches(5), Inches(0.4),
                "汇报人：张工  |  2025年12月", font_size=Pt(14), color=GRAY)
    # Bottom note
    add_textbox(slide, MARGIN, SLIDE_H - Inches(0.8), Inches(8), Inches(0.35),
                "CONFIDENTIAL · 宁德时代新能源科技股份有限公司", font_size=Pt(10), color=GRAY)


# ═══════════════════════════════════════════════
#  SLIDE 2 - TABLE OF CONTENTS
# ═══════════════════════════════════════════════
def make_toc():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_top_accent(slide)
    add_textbox(slide, MARGIN, Inches(0.4), Inches(10), Inches(0.7),
                "目  录", font_size=Pt(32), bold=True, color=DARK)

    items = [
        ("01", "工作概述", "岗位职责与团队背景"),
        ("02", "核心项目", "Eswin 售前储能管理系统"),
        ("03", "完成状况", "交付指标与数据"),
        ("04", "项目亮点", "技术突破与业务价值"),
        ("05", "不足与反思", "问题分析与改进方向"),
        ("06", "明年计划", "2026 年目标与规划"),
    ]

    start_y = Inches(1.6)
    for i, (num, title, desc) in enumerate(items):
        y = start_y + i * Inches(0.85)
        # Number circle
        add_bullet_circle(slide, MARGIN + Inches(0.1), y + Inches(0.02),
                          Inches(0.42), color=BLUE)
        add_textbox(slide, MARGIN + Inches(0.16), y + Inches(0.04),
                    Inches(0.35), Inches(0.35), num,
                    font_size=Pt(14), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        # Title
        add_textbox(slide, MARGIN + Inches(0.8), y, Inches(4), Inches(0.4),
                    title, font_size=Pt(20), bold=True, color=DARK)
        # Description
        add_textbox(slide, MARGIN + Inches(0.8), y + Inches(0.38), Inches(6), Inches(0.35),
                    desc, font_size=Pt(12), color=GRAY)

    add_page_number(slide, 2)


# ═══════════════════════════════════════════════
#  SLIDE 3 - 岗位概述
# ═══════════════════════════════════════════════
def make_role_overview():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_slide_title(slide, "岗位概述", "Role Overview")

    cards = [
        ("职位", ["前端开发工程师（高级）", "宁德时代 · 储能事业部"], []),
        ("主要职责", [
            "负责储能业务线 B 端管理系统前端架构与开发",
            "构建可复用业务组件库，提升团队开发效率",
            "前端工程化建设：CI/CD、代码规范、自动化测试",
            "与产品、后端团队协作，推进项目高质量交付",
        ], []),
        ("技术领域", ["React / TypeScript / Ant Design Pro", "中后台系统 · 表单引擎 · 数据可视化"], []),
    ]

    x_start = MARGIN
    card_w = Inches(3.6)
    card_h = Inches(2.8)
    gap = Inches(0.3)

    for i, (title, lines, _) in enumerate(cards):
        x = x_start + i * (card_w + gap)
        add_card(slide, x, Inches(1.5), card_w, card_h)

        # Card header with blue left accent
        accent_mini = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.5),
                                              Inches(0.08), card_h)
        accent_mini.fill.solid()
        accent_mini.fill.fore_color.rgb = BLUE
        accent_mini.line.fill.background()

        add_textbox(slide, x + Inches(0.25), Inches(1.65), card_w - Inches(0.4), Inches(0.4),
                    title, font_size=Pt(18), bold=True, color=DARK)
        add_multiline_textbox(slide, x + Inches(0.25), Inches(2.15),
                              card_w - Inches(0.4), card_h - Inches(0.8),
                              lines, font_size=Pt(13), color=GRAY, line_spacing=Pt(20))

    add_page_number(slide, 3)


# ═══════════════════════════════════════════════
#  SLIDE 4 - 核心项目
# ═══════════════════════════════════════════════
def make_core_project():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_slide_title(slide, "核心项目", "Core Project  ·  Eswin 售前储能管理系统")

    # Left: Project intro card
    left_w = Inches(5.2)
    add_card(slide, MARGIN, Inches(1.5), left_w, Inches(5.0))

    intro_lines = [
        "Eswin 售前储能管理系统是面向宁德时代售前团队的",
        "一站式储能方案配置与管理平台。系统支持从客户需求",
        "采集、储能方案智能配置、技术提案自动生成，到报价",
        "辅助决策的完整售前业务闭环。",
        "",
        "作为项目核心前端开发，我负责整体前端架构设计、",
        "复杂表单引擎开发、组件库建设及前端工程化落地。",
        "",
        "项目周期：2024 Q4 – 2025 Q3（已交付）",
        "团队规模：前端 3 人，整体 15 人",
    ]
    add_multiline_textbox(slide, MARGIN + Inches(0.25), Inches(1.7),
                          left_w - Inches(0.5), Inches(4.5),
                          intro_lines, font_size=Pt(14), color=DARK, line_spacing=Pt(22))

    # Right: Key metrics cards
    right_x = MARGIN + left_w + Inches(0.4)
    metric_w = Inches(2.6)
    metric_h = Inches(1.5)
    gap = Inches(0.3)

    metrics = [
        ("6 个", "核心功能模块", BLUE),
        ("12 次", "版本迭代", GREEN),
        ("~8 万行", "前端代码量", ORANGE),
        ("98%+", "Bug 修复率", GREEN),
        ("0 事故", "线上运行", GREEN),
        ("V1.0 如期", "交付里程碑", BLUE),
    ]

    positions = [
        (0, 0), (1, 0), (0, 1), (1, 1), (0, 2), (1, 2)
    ]
    for idx, (col, row) in enumerate(positions):
        if idx >= len(metrics):
            break
        num, lab, clr = metrics[idx]
        x = right_x + col * (metric_w + gap)
        y = Inches(1.5) + row * (metric_h + gap)
        add_metric_card(slide, x, y, metric_w, metric_h, num, lab, num_color=clr)

    add_page_number(slide, 4)


# ═══════════════════════════════════════════════
#  SLIDE 5 - 技术架构
# ═══════════════════════════════════════════════
def make_tech_stack():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_slide_title(slide, "技术架构", "Technical Architecture")

    categories = [
        ("核心框架", ["React 18 + TypeScript", "Ant Design Pro v5"]),
        ("状态管理", ["Zustand", "React Query (TanStack)"]),
        ("构建工具", ["Vite 5", "pnpm Workspace (Monorepo)"]),
        ("工程化", ["ESLint + Prettier + Stylelint", "Husky + lint-staged", "GitLab CI/CD"]),
        ("质量保障", ["Jest + Testing Library (75%+)", "Playwright E2E", "Storybook 组件文档"]),
    ]

    y_start = Inches(1.6)
    card_w = Inches(5.5)
    card_h = Inches(1.8)
    right_x = MARGIN + card_w + Inches(0.4)

    for i, (cat, items) in enumerate(categories):
        col = i % 2
        row = i // 2
        x = MARGIN if col == 0 else right_x
        y = y_start + row * Inches(2.0)

        # Card background
        bg = add_card(slide, x, y, card_w, card_h, fill_color=CARD_BG)

        # Blue dot indicator
        add_bullet_circle(slide, x + Inches(0.2), y + Inches(0.25),
                          Inches(0.18), color=BLUE)
        add_textbox(slide, x + Inches(0.55), y + Inches(0.2),
                    card_w - Inches(0.8), Inches(0.4),
                    cat, font_size=Pt(18), bold=True, color=DARK)

        # Items
        add_multiline_textbox(slide, x + Inches(0.55), y + Inches(0.7),
                              card_w - Inches(0.8), card_h - Inches(0.9),
                              [f"▸  {item}" for item in items],
                              font_size=Pt(13), color=GRAY, line_spacing=Pt(18))

    add_page_number(slide, 5)


# ═══════════════════════════════════════════════
#  SLIDE 6 - 功能模块交付
# ═══════════════════════════════════════════════
def make_modules():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_slide_title(slide, "功能模块交付", "Delivered Feature Modules")

    modules = [
        ("储能方案配置引擎", "JSON Schema 驱动的复杂表单系统\n支持动态字段、条件渲染、交叉校验", "独立负责前端架构"),
        ("技术提案生成模块", "模板引擎 + Word / PDF 导出\n支持复杂嵌套表格与图表", "核心开发 + 导出引擎"),
        ("客户需求管理系统", "全生命周期需求追踪\n工作流引擎 + 协作通知", "功能模块负责人"),
        ("产品数据库与选型", "储能产品参数库 + 智能推荐\n性能优化：万级数据流畅交互", "主导优化与组件抽象"),
        ("权限与角色管理", "RBAC 权限体系 + 动态路由\n按钮级权限控制", "架构设计与实现"),
        ("数据看板与分析", "ECharts 可视化 + 实时数据\n多维度业务指标展示", "独立开发交付"),
    ]

    y_start = Inches(1.4)
    card_w = Inches(3.7)
    card_h = Inches(2.6)
    gap_x = Inches(0.3)
    gap_y = Inches(0.3)

    for i, (name, desc, contrib) in enumerate(modules):
        col = i % 3
        row = i // 3
        x = MARGIN + col * (card_w + gap_x)
        y = y_start + row * (card_h + gap_y)

        add_card(slide, x, y, card_w, card_h)

        # Module number
        add_textbox(slide, x + Inches(0.2), y + Inches(0.15),
                    Inches(0.5), Inches(0.35),
                    f"0{i+1}", font_size=Pt(20), bold=True, color=BLUE)

        # Module name
        add_textbox(slide, x + Inches(0.2), y + Inches(0.55),
                    card_w - Inches(0.4), Inches(0.4),
                    name, font_size=Pt(15), bold=True, color=DARK)

        # Description
        add_multiline_textbox(slide, x + Inches(0.2), y + Inches(1.0),
                              card_w - Inches(0.4), Inches(1.0),
                              desc.split('\n'),
                              font_size=Pt(11), color=GRAY, line_spacing=Pt(16))

        # Contribution tag
        tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     x + Inches(0.2), y + card_h - Inches(0.5),
                                     Inches(2.0), Inches(0.3))
        tag.fill.solid()
        tag.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xFE)
        tag.line.fill.background()
        add_textbox(slide, x + Inches(0.3), y + card_h - Inches(0.48),
                    Inches(1.8), Inches(0.28),
                    contrib, font_size=Pt(9), color=BLUE, alignment=PP_ALIGN.CENTER)

    add_page_number(slide, 6)


# ═══════════════════════════════════════════════
#  SLIDE 7 - 项目完成状况
# ═══════════════════════════════════════════════
def make_project_status():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_slide_title(slide, "项目完成状况", "Project Delivery Status")

    # Big metrics row
    big_metrics = [
        ("100%", "整体进度\n按时完成全部里程碑", GREEN),
        ("6", "核心模块\n已交付上线", BLUE),
        ("12", "正式版本\n迭代发布", ORANGE),
    ]

    metric_w = Inches(3.5)
    metric_h = Inches(2.0)
    gap = Inches(0.4)
    total_w = len(big_metrics) * metric_w + (len(big_metrics) - 1) * gap
    x_start = (SLIDE_W - total_w) // 2

    for i, (num, lab, clr) in enumerate(big_metrics):
        x = x_start + i * (metric_w + gap)
        add_card(slide, x, Inches(1.5), metric_w, metric_h)
        add_textbox(slide, x, Inches(1.65), metric_w, Inches(0.7),
                    num, font_size=Pt(36), color=clr, bold=True,
                    alignment=PP_ALIGN.CENTER)
        lines = lab.split('\n')
        add_multiline_textbox(slide, x, Inches(2.35), metric_w, Inches(0.8),
                              lines, font_size=Pt(13), color=GRAY,
                              line_spacing=Pt(18))
        # Align center for these textboxes
        for para in slide.shapes:
            if hasattr(para, 'text_frame'):
                for p in para.text_frame.paragraphs:
                    if lab.split('\n')[0] in (p.text or ''):
                        p.alignment = PP_ALIGN.CENTER

    # Bottom small metrics
    small_metrics = [
        ("98%+", "P0/P1 Bug\n48h 内关闭"),
        ("75%+", "单元测试\n覆盖率"),
        ("92分", "Lighthouse\n性能评分"),
        ("40%", "打包体积\n缩减"),
    ]

    small_w = Inches(2.6)
    small_h = Inches(1.6)
    total_small = len(small_metrics) * small_w + (len(small_metrics) - 1) * gap
    xs = (SLIDE_W - total_small) // 2

    for i, (num, lab) in enumerate(small_metrics):
        x = xs + i * (small_w + gap)
        add_card(slide, x, Inches(3.9), small_w, small_h)
        add_textbox(slide, x, Inches(4.0), small_w, Inches(0.55),
                    num, font_size=Pt(24), color=DARK, bold=True,
                    alignment=PP_ALIGN.CENTER)
        lines = lab.split('\n')
        add_multiline_textbox(slide, x, Inches(4.55), small_w, Inches(0.7),
                              lines, font_size=Pt(11), color=GRAY,
                              line_spacing=Pt(16))

    # Timeline note
    add_textbox(slide, MARGIN, Inches(5.8), Inches(10), Inches(0.4),
                "▸  V1.0 提前 2 周交付，支撑重要客户技术评审  |  V1.5 性能优化专项完成，系统响应提升 60%",
                font_size=Pt(12), color=BLUE)

    add_page_number(slide, 7)


# ═══════════════════════════════════════════════
#  SLIDE 8 - 项目亮点 (I)
# ═══════════════════════════════════════════════
def make_highlights_1():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_slide_title(slide, "项目亮点（一）", "Technical Highlights")

    highlights = [
        ("复杂表单引擎", [
            "设计声明式 JSON Schema 驱动的表单引擎，将配置器开发效率提升 60%",
            "支持动态字段渲染、条件显隐、跨字段联动校验、自定义组件注册",
            "已应用于全部 6 个配置类模块，零重复代码",
        ]),
        ("共享组件库", [
            "构建 30+ 业务组件（储能领域定制），被 2 个兄弟项目复用",
            "配套 Storybook 文档 + 80% 单元测试覆盖",
            "统一视觉规范，消除各模块 UI 不一致问题",
        ]),
        ("提案导出引擎", [
            "模板驱动的 Word / PDF 导出系统，支持复杂嵌套表格与 ECharts 图表",
            "服务端渲染兜底方案，大文件导出稳定性 99.9%",
            "客户侧提案生成从 2 天缩短至 30 分钟",
        ]),
    ]

    y_start = Inches(1.5)
    card_w = Inches(10.8)
    card_h = Inches(1.7)
    gap = Inches(0.25)

    for i, (title, lines) in enumerate(highlights):
        y = y_start + i * (card_h + gap)
        add_card(slide, MARGIN, y, card_w, card_h, fill_color=CARD_BG)

        # Left accent strip
        strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, y,
                                        Inches(0.08), card_h)
        strip.fill.solid()
        strip.fill.fore_color.rgb = BLUE
        strip.line.fill.background()

        # Title
        add_textbox(slide, MARGIN + Inches(0.3), y + Inches(0.15),
                    Inches(3), Inches(0.35),
                    title, font_size=Pt(18), bold=True, color=DARK)

        # Content
        add_multiline_textbox(slide, MARGIN + Inches(0.3), y + Inches(0.55),
                              card_w - Inches(0.5), Inches(1.0),
                              [f"•  {l}" for l in lines],
                              font_size=Pt(13), color=GRAY, line_spacing=Pt(18))

    add_page_number(slide, 8)


# ═══════════════════════════════════════════════
#  SLIDE 9 - 项目亮点 (II)
# ═══════════════════════════════════════════════
def make_highlights_2():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_slide_title(slide, "项目亮点（二）", "Engineering & Performance")

    # Left card: Performance
    left_w = Inches(5.2)
    add_card(slide, MARGIN, Inches(1.5), left_w, Inches(4.5))

    add_textbox(slide, MARGIN + Inches(0.25), Inches(1.65),
                left_w - Inches(0.4), Inches(0.4),
                "⚡  性能优化", font_size=Pt(20), bold=True, color=DARK)
    perf_lines = [
        "Lighthouse 评分：45 → 92（+104%）",
        "首屏加载：3.2s → 1.1s（-66%）",
        "打包体积：4.8MB → 2.9MB（-40%）",
        "",
        "关键措施：",
        "• 路由级代码分割 + 组件懒加载",
        "• 虚拟滚动处理万级表格数据",
        "• 图片 WebP 自适应 + 懒加载",
        "• Tree Shaking + CDN 外链依赖",
    ]
    add_multiline_textbox(slide, MARGIN + Inches(0.25), Inches(2.15),
                          left_w - Inches(0.4), Inches(3.5),
                          perf_lines, font_size=Pt(13), color=GRAY,
                          line_spacing=Pt(18))

    # Right card: Engineering
    right_x = MARGIN + left_w + Inches(0.4)
    right_w = Inches(5.2)
    add_card(slide, right_x, Inches(1.5), right_w, Inches(4.5))

    add_textbox(slide, right_x + Inches(0.25), Inches(1.65),
                right_w - Inches(0.4), Inches(0.4),
                "🔧  工程化建设", font_size=Pt(20), bold=True, color=DARK)

    eng_lines = [
        "CI/CD 流水线（GitLab CI）",
        "• 自动 Lint → Test → Build → Deploy",
        "• MR 自动部署预览环境，提升 Review 效率",
        "",
        "代码规范体系",
        "• ESLint + Prettier + Stylelint 统一规范",
        "• Husky pre-commit 拦截不规范代码",
        "• 自动生成 Changelog + 版本管理",
        "",
        "成果：代码 Review 覆盖率 100%",
        "线上事故率 0，交付质量显著提升",
    ]
    add_multiline_textbox(slide, right_x + Inches(0.25), Inches(2.15),
                          right_w - Inches(0.4), Inches(3.5),
                          eng_lines, font_size=Pt(13), color=GRAY,
                          line_spacing=Pt(18))

    add_page_number(slide, 9)


# ═══════════════════════════════════════════════
#  SLIDE 10 - 不足与反思
# ═══════════════════════════════════════════════
def make_shortcomings():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_slide_title(slide, "不足与反思", "Shortcomings & Reflections")

    items = [
        ("自动化测试不足",
         ["E2E 测试覆盖偏低，核心流程缺乏回归保障",
          "已在 Q4 引入 Playwright，计划明年补齐关键路径"]),
        ("技术债务积累",
         ["部分早期模块存在代码耦合，重构优先级待排期",
          "表单引擎 TypeScript 类型推导仍有优化空间"]),
        ("文档沉淀滞后",
         ["组件库文档更新不及时，新人上手成本偏高",
          "后续将文档更新纳入 Definition of Done"]),
        ("跨团队协作效率",
         ["后端接口文档不规范导致联调返工",
          "已推动团队引入 Swagger / OpenAPI 规范"]),
    ]

    y_start = Inches(1.5)
    card_w = Inches(5.3)
    card_h = Inches(2.2)
    gap = Inches(0.3)

    for i, (title, lines) in enumerate(items):
        col = i % 2
        row = i // 2
        x = MARGIN + col * (card_w + gap)
        y = y_start + row * (card_h + gap)

        # Card with warm/red accent
        add_card(slide, x, y, card_w, card_h)

        # Top accent strip
        strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y,
                                        card_w, Inches(0.05))
        strip.fill.solid()
        strip.fill.fore_color.rgb = ORANGE
        strip.line.fill.background()

        # Number
        add_textbox(slide, x + Inches(0.2), y + Inches(0.2),
                    Inches(0.3), Inches(0.3),
                    f"0{i+1}", font_size=Pt(14), bold=True, color=ORANGE)

        # Title
        add_textbox(slide, x + Inches(0.55), y + Inches(0.2),
                    card_w - Inches(0.7), Inches(0.35),
                    title, font_size=Pt(16), bold=True, color=DARK)

        # Content
        add_multiline_textbox(slide, x + Inches(0.2), y + Inches(0.65),
                              card_w - Inches(0.4), Inches(1.3),
                              [f"•  {l}" for l in lines],
                              font_size=Pt(12), color=GRAY, line_spacing=Pt(18))

    add_page_number(slide, 10)


# ═══════════════════════════════════════════════
#  SLIDE 11 - 明年计划
# ═══════════════════════════════════════════════
def make_next_year_plan():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_slide_title(slide, "2026 年工作计划", "Next Year Plan")

    plans = [
        ("Eswin V2.0 架构升级", [
            "微前端改造（Module Federation），支持多租户",
            "接入 AI 辅助配置推荐能力",
            "移动端适配（储能现场勘查场景）",
        ], BLUE),
        ("技术深度提升", [
            "深入 Next.js SSR/SSG 服务端渲染方案",
            "探索 WebAssembly 在计算密集型场景的应用",
            "关注 React Server Components 生态演进",
        ], GREEN),
        ("质量体系建设", [
            "E2E 测试覆盖核心业务流程 100%",
            "单元测试覆盖率提升至 85%+",
            "接入前端性能监控（Sentry + 自研）",
        ], ORANGE),
        ("团队贡献与成长", [
            "每月 1 次前端技术分享，输出 2 篇技术博客",
            "完善脚手架工具，降低新项目启动成本",
            "参与面试与新人导师计划",
        ], BLUE),
    ]

    y_start = Inches(1.4)
    card_h = Inches(1.3)
    card_w = Inches(10.8)
    gap = Inches(0.2)

    for i, (title, lines, clr) in enumerate(plans):
        y = y_start + i * (card_h + gap)
        add_card(slide, MARGIN, y, card_w, card_h, fill_color=CARD_BG)

        # Left number circle
        add_bullet_circle(slide, MARGIN + Inches(0.15), y + Inches(0.35),
                          Inches(0.5), color=clr)
        add_textbox(slide, MARGIN + Inches(0.2), y + Inches(0.37),
                    Inches(0.4), Inches(0.35),
                    f"0{i+1}", font_size=Pt(14), color=WHITE, bold=True,
                    alignment=PP_ALIGN.CENTER)

        # Title
        add_textbox(slide, MARGIN + Inches(0.85), y + Inches(0.1),
                    Inches(4), Inches(0.35),
                    title, font_size=Pt(17), bold=True, color=DARK)

        # Items inline
        item_text = "    ".join([f"▸ {l}" for l in lines])
        add_textbox(slide, MARGIN + Inches(0.85), y + Inches(0.5),
                    card_w - Inches(1.1), Inches(0.6),
                    item_text, font_size=Pt(11), color=GRAY)

    add_page_number(slide, 11)


# ═══════════════════════════════════════════════
#  SLIDE 12 - 结语
# ═══════════════════════════════════════════════
def make_closing():
    slide = prs.slides.add_slide(BLANK_LAYOUT)

    # Subtle background circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     SLIDE_W // 2 - Inches(2.5), Inches(1.5),
                                     Inches(5), Inches(5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xFE)
    circle.line.fill.background()

    # Main text
    add_textbox(slide, Inches(2), Inches(2.5), Inches(9), Inches(1.0),
                "感谢团队的信任与支持", font_size=Pt(36), bold=True,
                color=DARK, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(2), Inches(3.6), Inches(9), Inches(0.6),
                "2026 年，继续深耕储能赛道，用技术驱动业务价值",
                font_size=Pt(18), color=GRAY, alignment=PP_ALIGN.CENTER)

    # Divider
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  SLIDE_W // 2 - Inches(1.5), Inches(4.5),
                                  Inches(3), Inches(0.03))
    div.fill.solid()
    div.fill.fore_color.rgb = BLUE
    div.line.fill.background()

    # Bottom text
    add_textbox(slide, Inches(2), Inches(5.0), Inches(9), Inches(0.5),
                "谢谢！", font_size=Pt(28), bold=True, color=BLUE,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(2), Inches(6.2), Inches(9), Inches(0.35),
                "宁德时代新能源科技股份有限公司  ·  储能事业部",
                font_size=Pt(12), color=GRAY, alignment=PP_ALIGN.CENTER)

    add_page_number(slide, 12)


# ═══════════════════════════════════════════════
#  BUILD ALL SLIDES
# ═══════════════════════════════════════════════
make_cover()
make_toc()
make_role_overview()
make_core_project()
make_tech_stack()
make_modules()
make_project_status()
make_highlights_1()
make_highlights_2()
make_shortcomings()
make_next_year_plan()
make_closing()

# Save
output_path = r"C:\Users\1\wechat-opencode\1\2025年度工作汇报.pptx"
prs.save(output_path)
print(f"[OK] PPT saved to: {output_path}")
print(f"     Slides: {len(prs.slides)}")
